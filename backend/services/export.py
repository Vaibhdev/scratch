import io
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import models

def create_docx(document_data: models.Document):
    doc = Document()
    doc.add_heading(document_data.project.title, 0)
    
    for section in document_data.sections:
        doc.add_heading(section.title, level=1)
        if section.content:
            doc.add_paragraph(section.content)
            
    file_stream = io.BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return file_stream

def create_pptx(document_data: models.Document):
    prs = Presentation()
    
    # Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = document_data.project.title
    subtitle.text = document_data.project.description or ""
    
    # Content Slides
    bullet_slide_layout = prs.slide_layouts[1]
    
    for section in document_data.sections:
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = section.title
        
        if section.content:
            tf = body_shape.text_frame
            tf.text = section.content
            
    file_stream = io.BytesIO()
    prs.save(file_stream)
    file_stream.seek(0)
    return file_stream
